import io
import time
import streamlit as st
import pandas as pd

# Parsers
from docx import Document           # from python-docx
from pptx import Presentation       # from python-pptx
import fitz                         # from pymupdf

st.set_page_config(page_title="M365 Accessibility Pre-Flight Check", layout="wide")
st.title("M365 Accessibility Pre-Flight Check")
st.write("Upload your documents (.docx, .pptx, .pdf) to check for basic accessibility issues.")

# -----------------------------
# DOCX checks
# -----------------------------
def check_docx_basic(file_obj):
    """Basic DOCX checks: headings present? ambiguous link text?"""
    issues = []
    doc = Document(file_obj)

    # Headings present?
    has_heading = any(p.style and str(p.style.name).startswith("Heading") for p in doc.paragraphs)
    if not has_heading:
        issues.append("No headings found (use Heading 1/2/3 styles for document structure).")

    # Ambiguous link text (very simple scan)
    for p in doc.paragraphs:
        for run in p.runs:
            t = (run.text or "").strip().lower()
            if t in {"here", "click here", "read more"}:
                issues.append("Avoid vague link text like 'click here'—use descriptive link text.")

    return issues


def check_docx_missing_alt_text(file_bytes: bytes):
    """
    Best-effort scan of inline images to see if docPr@descr (alt text) is present.
    Returns a list of issue strings (or an empty list).
    """
    issues = []
    doc = Document(io.BytesIO(file_bytes))

    total_imgs = 0
    missing = 0
    # Try the python-docx inline_shapes path first
    try:
        for ish in doc.inline_shapes:
            total_imgs += 1
            # Access the underlying XML (docx uses lxml under the hood)
            descr = ish._inline.docPr.get("descr")  # type: ignore[attr-defined]
            if not (descr and descr.strip()):
                missing += 1
    except Exception:
        # Fallback: count related images (cannot read alt reliably this way)
        img_rels = [r for r in doc.part.rels.values() if "image" in r.reltype]
        total_imgs = len(img_rels)
        # When we can't inspect alt text, warn the user to verify
        if total_imgs:
            issues.append(
                f"Found {total_imgs} image(s). Verify Alt Text for each (Right-click image → View Alt Text)."
            )
            return issues

    if total_imgs > 0:
        if missing > 0:
            issues.append(
                f"{missing} of {total_imgs} image(s) appear to lack alt text. "
                "Add concise, purpose-focused descriptions (≤125 characters)."
            )
        else:
            issues.append(f"All {total_imgs} image(s) have alt text. ✅")

    return issues


# -----------------------------
# PPTX checks
# -----------------------------
def check_pptx(file_obj):
    """Basic PPTX checks: slide titles present? remind about alt text & contrast."""
    issues = []
    prs = Presentation(file_obj)

    # Slide titles as heading proxy
    missing_titles = []
    for i, slide in enumerate(prs.slides, start=1):
        has_title = getattr(slide.shapes, "title", None) is not None
        any_text = any(getattr(s, "has_text_frame", False) and s.text_frame and s.text_frame.text.strip()
                       for s in slide.shapes)
        if not has_title and not any_text:
            missing_titles.append(i)
    if missing_titles:
        issues.append(f"Slides without a clear title/text: {missing_titles}. Use a Title layout or add a heading.")

    # Image reminder (python-pptx does not expose alt text reliably for all shapes)
    img_count = 0
    for slide in prs.slides:
        for s in slide.shapes:
            if getattr(s, "image", None) or getattr(s, "shape_type", None) in {13, 14}:
                img_count += 1
    if img_count:
        issues.append(
            f"Found {img_count} image(s). Verify Alt Text for each (Format Picture → Alt Text)."
        )

    # Contrast placeholder (add pixel-based checks later if you want)
    issues.append("Contrast not evaluated in MVP. Aim for WCAG contrast ratio ≥ 4.5:1 for normal text.")

    return issues


# -----------------------------
# PDF checks
# -----------------------------
def check_pdf(file_obj):
    """Basic PDF checks: tagged structure, rough heading proxy, link-text reminder."""
    issues = []
    file_bytes = file_obj.read()
    pdf = fitz.open(stream=file_bytes, filetype="pdf")

    # Tagged PDF?
    if not pdf.is_tagged:
        issues.append("PDF is not tagged (no accessibility structure). Export with 'Create tagged PDF' enabled.")

    # Heuristic: look for large text spans as rough heading proxy
    big, total = 0, 0
    for page in pdf:
        blocks = page.get_text("dict").get("blocks", [])
        for b in blocks:
            for l in b.get("lines", []):
                for s in l.get("spans", []):
                    text = (s.get("text") or "").strip()
                    if not text:
                        continue
                    total += 1
                    if s.get("size", 0) >= 16:
                        big += 1
    if total and (big / total) < 0.02:
        issues.append("Few/no large text spans detected; headings may be missing. Use heading styles in the source doc.")

    # General reminder on link text (PDF link text extraction is unreliable)
    issues.append("Review links for descriptive text (avoid 'click here'). Edit in source, then re-export.")
    return issues


# -----------------------------
# File uploader & processing
# -----------------------------
uploaded_files = st.file_uploader(
    "Upload files", type=["docx", "pptx", "pdf"], accept_multiple_files=True
)

results = []
if uploaded_files:
    for file in uploaded_files:
        file_name = file.name
        issues_list = []

        if file_name.lower().endswith(".docx"):
            # 1) Basic DOCX checks
            issues_basic = check_docx_basic(file)

            # 2) DOCX missing-alt-text check uses raw bytes
            file.seek(0)
            docx_bytes = file.read()
            issues_alt = check_docx_missing_alt_text(docx_bytes)

            # Merge results
            issues_list = (issues_basic or []) + (issues_alt or [])

            # Reset pointer in case you read again later
            file.seek(0)

        elif file_name.lower().endswith(".pptx"):
            issues_list = check_pptx(file)

        elif file_name.lower().endswith(".pdf"):
            file.seek(0)
            issues_list = check_pdf(file)
            file.seek(0)

        else:
            issues_list = ["Unsupported file type."]

        # Ensure we always produce at least one message
        if not issues_list:
            issues_list = ["No major issues found."]

        results.append({"File": file_name, "Issues": "; ".join(issues_list)})

# -----------------------------
# Results table + downloads
# -----------------------------
if results:
    df = pd.DataFrame(results)
    st.subheader("Accessibility Findings")
    st.dataframe(df, use_container_width=True)

    # Downloadable CSV report
    csv = df.to_csv(index=False).encode("utf-8")
    st.download_button("Download CSV Report", data=csv, file_name="accessibility_report.csv", mime="text/csv")

    # Downloadable HTML report
    html = df.to_html(index=False)
    st.download_button("Download HTML Report", data=html, file_name="accessibility_report.html", mime="text/html")
